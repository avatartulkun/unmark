from __future__ import annotations

"""去掉 PPTX 里跨页固定的角标水印（NotebookLM 导出等）。全程离线。

和 PDF 那条路同一个判据——**内容在变、它不变**。但 PPTX 有两种长相，走两条不同的路：

  1. **水印是个形状对象**（文本框、图片、图形）→ 直接删掉。见 `clean_pptx`
  2. **每页就是一整张图，水印烧在像素里** → 和 PDF 完全一样的处理，
     只是图从 zip 里掏、修完再塞回去。见 `clean_pptx_bitmaps`

第二种是真实导出件里更常见的那一种（实测 NotebookLM 导出的绘本 pptx 就是），
所以 `clean_pptx` 找不到对象时会自动转到它。

先说第一种，它更干净：

PPTX 里的水印不是像素，是一个形状对象。所以这里不涂、不填、不重编码，
**直接把那个对象删掉**。后果是：

  · 无损。其余每一个字、每一张图，二进制原样保留
  · 文件不会变大（PDF 那条路因为要无损重编码位图，输出常常涨到几倍）
  · 没有「修复得像不像」这种问题，也就没有磨砂衬底那一堆麻烦

母版和版式上的水印也算数：版式上的一个形状会出现在所有用它的页面上，
所以按「它实际覆盖到多少页」来算证据，和逐页放置的水印同一个标准。

安全保证与 PDF 那条路一致：原文件永不改写，结果写到新文件；写出后重新打开，
逐页核对「该留的一个没少、该删的一个没剩」。
"""

import hashlib
import re
import shutil
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Optional, Sequence

import cv2
import numpy as np
from pptx import Presentation

from unwatermark import CleanOptions, detect_in_frames, repair_frame

ProgressCallback = Optional[Callable[[int, int, str], None]]

_WHITESPACE = re.compile(r"\s+")
_XML_ID = re.compile(rb'\sid="\d+"')
_XML_NAME = re.compile(rb'\sname="[^"]*"')


@dataclass(frozen=True)
class PptxOptions:
    """检测参数；默认值针对导出工具加在角落的那种角标。"""

    vote: float = 0.85
    """一个形状要覆盖到多大比例的页面，才算「跨页固定」。

    和 PDF 那条路取同一个值。留出余量是因为封面、封底常常用不同的版式，
    水印在那几页可能确实没有。
    """
    min_slides: int = 3
    """少于这么多页就不做自动判定——两页之间「都有」不足以说明什么。"""
    max_area: float = 0.08
    """自动模式下，一个候选最多能占页面多大面积。角标是小东西。"""
    margin_w: float = 0.25
    margin_h: float = 0.15
    """页面中央「正文区」的边距比例。自动模式只动完全落在这块之外的形状。

    水印待在页边，正文待在中间。压到正文区的东西就算跨页固定，也更可能是
    版式元素（标题栏、装饰块），不该替用户做主删掉。
    """
    outer_only: bool = True
    """自动模式是否只考虑落在正文区之外的形状。关掉就只剩「跨页固定 + 够小」。"""


@dataclass(frozen=True)
class Mark:
    """一处跨页固定的元素。`key` 用来在调用方和引擎之间指认它。"""

    key: str
    kind: str
    """text / picture / shape。"""
    label: str
    """给人看的名字：文字内容，或「图片 12.4 KB」。"""
    home: str
    """它长在哪儿：slide（逐页放置）/ layout（版式）/ master（母版）。"""
    hits: int
    """实际覆盖到多少页。"""
    total: int
    rect: tuple[int, int, int, int]
    """位置和尺寸，EMU（1 英寸 = 914400）。"""
    area_percent: float
    outer: bool
    """是否完全落在正文区之外。"""
    eligible: bool
    """自动模式会不会动它。"""
    reason: str
    """不合格时说明原因，合格时为空。"""


@dataclass(frozen=True)
class PptxResult:
    success: bool
    message: str
    slides: int = 0
    removed: int = 0
    marks: list[Mark] = field(default_factory=list)
    removed_keys: list[str] = field(default_factory=list)
    output_path: Optional[Path] = None
    mode: str = "objects"
    """objects：删掉了形状对象。bitmaps：整页位图型，走的是 PDF 那一套。"""
    strategy: str = ""
    box: Optional[tuple[int, int, int, int]] = None
    box_ratio: Optional[tuple[float, float, float, float]] = None
    area_percent: float = 0.0
    pages_processed: int = 0

    def as_dict(self) -> dict:
        return {
            "success": self.success,
            "message": self.message,
            "slides": self.slides,
            "removed": self.removed,
            "removed_keys": list(self.removed_keys),
            "mode": self.mode,
            "strategy": self.strategy,
            "box": list(self.box) if self.box else None,
            "box_ratio": list(self.box_ratio) if self.box_ratio else None,
            "area_percent": self.area_percent,
            "pages_processed": self.pages_processed,
            "marks": [
                {
                    "key": m.key, "kind": m.kind, "label": m.label, "home": m.home,
                    "hits": m.hits, "total": m.total, "area_percent": m.area_percent,
                    "outer": m.outer, "eligible": m.eligible, "reason": m.reason,
                }
                for m in self.marks
            ],
            "output": str(self.output_path) if self.output_path else None,
        }


# ---------------------------------------------------------------- 形状指纹

def _text_of(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    try:
        return _WHITESPACE.sub(" ", shape.text_frame.text or "").strip()
    except Exception:
        return ""


def _content_key(shape) -> tuple[str, str]:
    """返回 (类别, 内容指纹)。指纹要跨页稳定，所以不能带对象 id。"""
    text = _text_of(shape)
    if text:
        return "text", hashlib.sha1(text.encode("utf-8")).hexdigest()[:16]
    try:                                       # 图片按像素内容认，改名换 id 都躲不掉
        blob = shape.image.blob
        return "picture", hashlib.sha1(blob).hexdigest()[:16]
    except Exception:
        pass
    try:
        # 兜底：拿形状自己的 XML 当指纹，但要抹掉逐页递增的 id 和自动生成的 name，
        # 否则同一个水印在每页的指纹都不一样，跨页交集永远为空。
        xml = shape._element.xml.encode("utf-8")
        xml = _XML_NAME.sub(b"", _XML_ID.sub(b"", xml))
        return "shape", hashlib.sha1(xml).hexdigest()[:16]
    except Exception:
        return "shape", "unknown"


def _rect_of(shape) -> Optional[tuple[int, int, int, int]]:
    """位置四元组，取整到百分之一英寸——同一个水印在各页的坐标本该完全相同，
    取整只是为了容忍导出工具偶尔差一两个 EMU。"""
    try:
        values = (shape.left, shape.top, shape.width, shape.height)
    except Exception:
        return None
    if any(v is None for v in values):
        return None                            # 继承来的占位符，没有自己的坐标
    grid = 9144                                # 1/100 英寸
    return tuple(int(round(int(v) / grid)) * grid for v in values)


def _signature(shape) -> Optional[tuple]:
    rect = _rect_of(shape)
    if rect is None:
        return None
    kind, content = _content_key(shape)
    return (kind, content, rect)


def _label_of(shape, kind: str) -> str:
    text = _text_of(shape)
    if text:
        return text if len(text) <= 60 else text[:57] + "…"
    if kind == "picture":
        try:
            return f"图片 {len(shape.image.blob) / 1024:.1f} KB"
        except Exception:
            return "图片"
    name = (getattr(shape, "name", "") or "").strip()
    return f"图形「{name}」" if name else "图形"


# ---------------------------------------------------------------- 扫描

def _visible_shapes(container) -> list:
    """容器里会真正画出来的顶层形状。

    版式和母版上的占位符是「提示框」，不填内容就不出现在页面上，所以排除；
    页面自己的占位符已经有内容，照常参与。
    """
    out = []
    for shape in container.shapes:
        try:
            if getattr(shape, "is_placeholder", False):
                continue
        except Exception:
            pass
        out.append(shape)
    return out


def _shows_master(part) -> bool:
    """版式/页面有没有关掉「显示母版图形」。"""
    element = getattr(part, "element", None)
    if element is None:
        return True
    return element.get("showMasterSp", "1") not in ("0", "false")


def scan_pptx(source: Path | str, options: PptxOptions = PptxOptions()) -> tuple[list[Mark], int, str]:
    """列出这份 PPTX 里所有跨页固定的元素。返回 (候选表, 页数, 出错原因)。

    候选表按「覆盖页数 → 面积从小到大」排，最像角标的排在前面。
    """
    presentation = Presentation(str(source))
    slides = list(presentation.slides)
    total = len(slides)
    if total == 0:
        return [], 0, "这份文件里没有幻灯片"

    width = int(presentation.slide_width or 0)
    height = int(presentation.slide_height or 0)
    if width <= 0 or height <= 0:
        return [], total, "取不到页面尺寸"

    # 每个指纹记：覆盖到哪些页、它长在哪儿、一个代表形状
    hits: dict[tuple, set[int]] = {}
    home: dict[tuple, str] = {}
    sample: dict[tuple, object] = {}

    def note(signature, index: int, where: str, shape) -> None:
        if signature is None:
            return
        hits.setdefault(signature, set()).add(index)
        # 逐页放置的优先级最低，母版最高——同一个指纹如果两处都有，
        # 删源头（母版/版式）才是一次性解决。
        rank = {"slide": 0, "layout": 1, "master": 2}
        if signature not in home or rank[where] > rank[home[signature]]:
            home[signature] = where
            sample[signature] = shape
        sample.setdefault(signature, shape)

    for index, slide in enumerate(slides):
        for shape in _visible_shapes(slide):
            note(_signature(shape), index, "slide", shape)
        layout = getattr(slide, "slide_layout", None)
        if layout is None:
            continue
        for shape in _visible_shapes(layout):
            note(_signature(shape), index, "layout", shape)
        if not (_shows_master(layout) and _shows_master(slide)):
            continue
        master = getattr(layout, "slide_master", None)
        if master is None:
            continue
        for shape in _visible_shapes(master):
            note(_signature(shape), index, "master", shape)

    # 正文区：中间那一块。水印待在页边，正文待在中间。
    body = (width * options.margin_w, height * options.margin_h,
            width * (1 - options.margin_w), height * (1 - options.margin_h))

    marks: list[Mark] = []
    for signature, pages in hits.items():
        count = len(pages)
        if count < max(2, int(total * options.vote + 0.5)):
            continue
        kind, content, rect = signature
        left, top, shape_w, shape_h = rect
        area = 100.0 * (shape_w * shape_h) / float(width * height)
        outer = not (left < body[2] and left + shape_w > body[0]
                     and top < body[3] and top + shape_h > body[1])

        reason = ""
        if total < options.min_slides:
            reason = f"只有 {total} 页，跨页比对不可靠"
        elif area > options.max_area * 100:
            reason = f"占了页面 {area:.1f}%，比角标大得多"
        elif options.outer_only and not outer:
            reason = "压在正文区里，更像版式元素"
        marks.append(Mark(
            key=f"{kind}:{content}:{left}:{top}",
            kind=kind,
            label=_label_of(sample[signature], kind),
            home=home[signature],
            hits=count,
            total=total,
            rect=rect,
            area_percent=area,
            outer=outer,
            eligible=not reason,
            reason=reason,
        ))

    marks.sort(key=lambda m: (-m.hits, m.area_percent))
    return marks, total, ""


# ---------------------------------------------------------------- 清理

def _drop(shape) -> None:
    parent = shape._element.getparent()
    if parent is not None:
        parent.remove(shape._element)


def _inventory(path: Path) -> tuple[int, list[set[tuple]], set[str]]:
    """把一份 PPTX 的内容摊平成可比对的形式，用来做写出后的复核。"""
    presentation = Presentation(str(path))
    per_slide: list[set[tuple]] = []
    texts: set[str] = set()
    for slide in presentation.slides:
        signatures = set()
        for shape in _visible_shapes(slide):
            signature = _signature(shape)
            if signature is not None:
                signatures.add(signature)
            text = _text_of(shape)
            if text:
                texts.add(text)
        per_slide.append(signatures)
    for master in presentation.slide_masters:
        for container in [master] + list(master.slide_layouts):
            for shape in _visible_shapes(container):
                text = _text_of(shape)
                if text:
                    texts.add(text)
    return len(per_slide), per_slide, texts


def clean_pptx(
    source: Path | str,
    destination: Path | str,
    options: PptxOptions = PptxOptions(),
    select: Optional[Sequence[str]] = None,
    progress: ProgressCallback = None,
    bitmap_options: Optional[CleanOptions] = None,
) -> PptxResult:
    """删掉跨页固定的角标，写出新文件。原文件永不改写。

    `select` 给定要删的 `Mark.key` 列表；不给就用自动模式——只删 `eligible` 的那些。

    找不到可删的对象时，自动转去试「整页位图型」——真实导出件里那种更常见，
    整页画面被压成一张位图、水印烧在像素里，PPTX 只是个壳。
    """
    source, destination = Path(source), Path(destination)
    if progress:
        progress(0, 1, "分析各页")
    marks, total, why = scan_pptx(source, options)
    if why:
        return PptxResult(False, why, slides=total)

    if select is None:
        wanted = {m.key for m in marks if m.eligible}
    else:
        wanted = set(select)
        unknown = wanted - {m.key for m in marks}
        if unknown:
            return PptxResult(False, f"指定的对象不在候选里：{sorted(unknown)[:3]}",
                              slides=total, marks=marks)
    if not wanted and select is None:
        # 没有可删的对象，不代表这份文件没水印——它可能整页就是一张图
        fallback = clean_pptx_bitmaps(source, destination, bitmap_options or CleanOptions(),
                                      progress)
        if fallback.success:
            return PptxResult(**{**fallback.__dict__, "marks": marks})
        hint = "；".join(f"「{m.label}」{m.reason}" for m in marks[:3] if m.reason)
        return PptxResult(
            False,
            "没找到跨页固定的角标" + (f"（最接近的：{hint}）" if hint else "")
            + f"；按整页位图处理也不行（{fallback.message}）",
            slides=total, marks=marks,
        )
    if not wanted:
        return PptxResult(False, "没有勾选任何对象", slides=total, marks=marks)

    chosen = {m.key: m for m in marks if m.key in wanted}
    before_digest = hashlib.sha1(source.read_bytes()).hexdigest()

    presentation = Presentation(str(source))
    slides = list(presentation.slides)
    removed = 0

    def sweep(container, where: str) -> None:
        nonlocal removed
        for shape in list(_visible_shapes(container)):
            signature = _signature(shape)
            if signature is None:
                continue
            kind, content, rect = signature
            key = f"{kind}:{content}:{rect[0]}:{rect[1]}"
            mark = chosen.get(key)
            if mark is not None and mark.home == where:
                _drop(shape)
                removed += 1

    for index, slide in enumerate(slides):
        if progress:
            progress(index + 1, len(slides), f"清理第 {index + 1} 页")
        sweep(slide, "slide")
    for master in presentation.slide_masters:
        sweep(master, "master")
        for layout in master.slide_layouts:
            sweep(layout, "layout")

    if removed == 0:
        return PptxResult(False, "没有对象被删除，原文件未改动", slides=total, marks=marks)

    destination.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(destination))

    if hashlib.sha1(source.read_bytes()).hexdigest() != before_digest:
        raise RuntimeError("安全检查失败：原文件被改动了。")

    if progress:
        progress(len(slides), len(slides), "复核")
    problems = _verify(source, destination, chosen.values(), options)
    if problems:
        destination.unlink(missing_ok=True)
        return PptxResult(False, "已中止：" + "；".join(problems), slides=total, marks=marks)

    where = "、".join(sorted({{"slide": "逐页", "layout": "版式", "master": "母版"}[m.home]
                             for m in chosen.values()}))
    listed = "、".join(f"「{m.label}」" for m in list(chosen.values())[:3])
    return PptxResult(
        True,
        f"删掉 {len(chosen)} 个跨页固定对象（{where}）：{listed}，共 {removed} 处，{total} 页已处理",
        slides=total, removed=removed, marks=marks,
        removed_keys=sorted(chosen), output_path=destination,
    )


def _verify(source: Path, destination: Path, chosen, options: PptxOptions) -> list[str]:
    """写出后重新打开核对：该留的一个没少，该删的一个没剩。

    删除是很容易做过头的操作——选错指纹就会把正文里的东西一起带走，而且不像
    涂像素那样一眼看得出来。所以这里不信内存里的状态，只信重新读回来的文件。
    """
    problems: list[str] = []
    src_count, src_slides, _src_texts = _inventory(source)
    out_count, out_slides, out_texts = _inventory(destination)
    if src_count != out_count:
        return [f"页数从 {src_count} 变成了 {out_count}"]

    targets = {m.key for m in chosen}
    for index, (before, after) in enumerate(zip(src_slides, out_slides)):
        expected = {s for s in before
                    if f"{s[0]}:{s[1]}:{s[2][0]}:{s[2][1]}" not in targets}
        if after != expected:
            lost = len(expected - after)
            extra = len(after - expected)
            problems.append(f"第 {index + 1} 页少了 {lost} 个、多了 {extra} 个对象")
            break

    for mark in chosen:
        if mark.kind == "text" and mark.label.rstrip("…") and not mark.label.endswith("…"):
            if mark.label in out_texts:
                problems.append(f"「{mark.label}」在清理后仍然存在")
                break
    return problems


# ---------------------------------------------------------------- 整页位图的 PPTX

@dataclass(frozen=True)
class BitmapDeck:
    """每页正好一张满版图的 PPTX。`parts` 是各页图片在 zip 里的路径。"""

    parts: list[str]
    height: int
    width: int


def _bitmap_layout(source: Path, coverage: float = 0.98) -> tuple[list[str], str]:
    """只看结构，不解码图片：每页是不是正好一张满版图？是的话返回各页图片的 zip 路径。

    和 `_bitmap_deck` 分开是有原因的：网页端每次预览都要问一次「这是不是整页位图型」，
    而解码十几张 PNG 要好几秒——预览接口被这个拖垮过。
    """
    presentation = Presentation(str(source))
    slides = list(presentation.slides)
    if len(slides) < 3:
        return [], f"只有 {len(slides)} 页，跨页比对不可靠"
    page_area = int(presentation.slide_width or 0) * int(presentation.slide_height or 0)
    if page_area <= 0:
        return [], "取不到页面尺寸"

    parts_out: list[str] = []
    for index, slide in enumerate(slides):
        pictures = [sh for sh in slide.shapes if sh.shape_type == 13]
        if len(pictures) != 1 or len(_visible_shapes(slide)) != 1:
            return [], f"第 {index + 1} 页不是「整页一张图」"
        picture = pictures[0]
        rect = _rect_of(picture)
        if rect is None:
            return [], f"第 {index + 1} 页的图片取不到位置"
        covered = rect[2] * rect[3]
        if covered < coverage * page_area:
            percent = 100.0 * covered / page_area
            return [], f"第 {index + 1} 页的图片只覆盖 {percent:.0f}%，不是满版"
        try:
            # partname 形如 /ppt/media/image1.png，去掉开头的斜杠就是 zip 里的路径
            related = slide.part.related_part(picture._element.blip_rId)
            parts_out.append(str(related.partname)[1:])
        except Exception:
            return [], f"第 {index + 1} 页取不到图片数据"
    return parts_out, ""


def _bitmap_deck(source: Path, coverage: float = 0.98) -> tuple[Optional[BitmapDeck], str]:
    """在结构判定之上再确认各页位图尺寸一致——跨页求交集的前提。"""
    parts, why = _bitmap_layout(source, coverage)
    if not parts:
        return None, why
    sizes = set()
    with zipfile.ZipFile(source) as archive:
        for name in parts:
            image = cv2.imdecode(np.frombuffer(archive.read(name), np.uint8), cv2.IMREAD_COLOR)
            if image is None:
                return None, f"{name} 解码失败"
            sizes.add(image.shape[:2])
    if len(sizes) != 1:
        return None, "各页位图尺寸不一致，跨页比对不成立"
    height, width = sizes.pop()
    return BitmapDeck(parts, height, width), ""


def bitmap_parts(source: Path | str) -> list[str]:
    """整页位图型 PPTX 各页图片在 zip 里的路径；不是这种就返回空表。

    网页端的「处理前 / 处理后」对比要靠它取图——这种 PPTX 每页正好一张满版图，
    所谓「渲染一页」就是把那张图读出来。
    """
    parts, _why = _bitmap_layout(Path(source))
    return parts


def clean_pptx_bitmaps(
    source: Path | str,
    destination: Path | str,
    options: CleanOptions = CleanOptions(),
    progress: ProgressCallback = None,
) -> PptxResult:
    """整页位图型 PPTX：把各页的图掏出来，按 PDF 那一套修好，再原位塞回去。

    检测和修复**完全复用** unwatermark 那边的代码——两种容器里的角标是同一个东西，
    判据也就该是同一份代码。这里只负责把图掏出来、塞回去。

    塞回去用的是「重打一遍 zip，只替换那几个图片条目」：其余每一个部件（版式、母版、
    关系表、主题、字体）都按原样逐字节拷过去，不经过 python-pptx 重新序列化——
    重新序列化会顺手改写一堆无关的 XML，diff 里根本看不出到底动了什么。
    """
    source, destination = Path(source), Path(destination)
    deck, why = _bitmap_deck(source)
    if deck is None:
        return PptxResult(False, why)

    if progress:
        progress(0, len(deck.parts), "分析各页")
    with zipfile.ZipFile(source) as archive:
        images = [cv2.cvtColor(cv2.imdecode(np.frombuffer(archive.read(name), np.uint8),
                                            cv2.IMREAD_COLOR), cv2.COLOR_BGR2RGB)
                  for name in deck.parts]

    def frames():
        for index, image in enumerate(images):
            if progress:
                progress(index + 1, len(images), f"分析第 {index + 1} 页")
            yield str(index), image

    detection, why = detect_in_frames(frames(), deck.height, deck.width, options)
    if detection is None:
        return PptxResult(False, f"没找到跨页固定的角标（{why}）", slides=len(images))

    x0, y0, x1, y1 = detection.box
    painted = detection.mask.copy()
    repaired: dict[str, bytes] = {}
    for index, (name, image) in enumerate(zip(deck.parts, images)):
        if progress:
            progress(index + 1, len(images), f"修复第 {index + 1} 页")
        fixed, used = repair_frame(image, detection.mask, options, detection.box)
        painted = np.maximum(painted, used)
        ok, buffer = cv2.imencode(Path(name).suffix or ".png",
                                  cv2.cvtColor(fixed, cv2.COLOR_RGB2BGR),
                                  [cv2.IMWRITE_PNG_COMPRESSION, options.png_compression])
        if not ok:
            return PptxResult(False, f"第 {index + 1} 页重新编码失败，原文件未改动",
                              slides=len(images))
        repaired[name] = buffer.tobytes()

    total = int((painted > 0).sum())
    area = 100.0 * total / float(deck.height * deck.width)
    if total > 4 * max(1, int((detection.mask > 0).sum())) and \
            total > max(options.panel_max_area, options.defrost_max_area) * deck.height * deck.width:
        return PptxResult(False, f"修复区异常膨胀到整页 {area:.2f}%，已中止", slides=len(images))

    destination.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(source) as src, \
            zipfile.ZipFile(destination, "w", zipfile.ZIP_DEFLATED) as out:
        for item in src.infolist():
            data = repaired.get(item.filename)
            if data is None:
                out.writestr(item, src.read(item.filename))   # 原样拷贝，连压缩方式都留着
            else:
                out.writestr(item.filename, data)

    problems = _verify_bitmaps(source, destination, deck, painted, options)
    if problems:
        destination.unlink(missing_ok=True)
        return PptxResult(False, "已中止：" + "；".join(problems), slides=len(images))

    return PptxResult(
        True,
        f"{detection.strategy}命中 x{x0}-{x1} y{y0}-{y1}，涂改 {area:.3f}%，"
        f"{len(images)} 页已处理（整页位图型 PPTX）",
        slides=len(images), removed=len(images), output_path=destination,
        mode="bitmaps", strategy=detection.strategy, box=detection.box,
        box_ratio=(x0 / deck.width, y0 / deck.height, x1 / deck.width, y1 / deck.height),
        area_percent=area, pages_processed=len(images),
    )


def _verify_bitmaps(source: Path, destination: Path, deck: BitmapDeck,
                    painted: np.ndarray, options: CleanOptions) -> list[str]:
    """写出后重新打开抽页复核：掩膜之外确实一个像素都没变，其余部件逐字节没动。"""
    problems: list[str] = []
    outside = painted == 0
    step = max(1, len(deck.parts) // max(1, options.verify_sample))
    with zipfile.ZipFile(source) as src, zipfile.ZipFile(destination) as out:
        src_names = {i.filename for i in src.infolist()}
        out_names = {i.filename for i in out.infolist()}
        if src_names != out_names:
            return [f"部件清单变了（少 {len(src_names - out_names)} 个、"
                    f"多 {len(out_names - src_names)} 个）"]
        touched = set(deck.parts)
        for name in sorted(src_names - touched):
            if src.read(name) != out.read(name):
                problems.append(f"不该改动的部件 {name} 变了")
                break
        for index in range(0, len(deck.parts), step):
            name = deck.parts[index]
            before = cv2.imdecode(np.frombuffer(src.read(name), np.uint8), cv2.IMREAD_COLOR)
            after = cv2.imdecode(np.frombuffer(out.read(name), np.uint8), cv2.IMREAD_COLOR)
            if after is None or after.shape != before.shape:
                problems.append(f"第 {index + 1} 页尺寸变了")
                break
            if not np.array_equal(before[outside], after[outside]):
                changed = int((before[outside] != after[outside]).any(axis=-1).sum())
                problems.append(f"第 {index + 1} 页 Mask 之外有 {changed} 个像素被改动")
                break
    return problems


def main(argv: Optional[Sequence[str]] = None) -> int:
    import argparse

    parser = argparse.ArgumentParser(description="去掉 PPTX 里跨页固定的角标水印")
    parser.add_argument("source")
    parser.add_argument("destination", nargs="?")
    parser.add_argument("--list", action="store_true", help="只列出候选，不改动")
    parser.add_argument("--select", action="append", default=None, help="指定要删的 key，可重复")
    args = parser.parse_args(argv)

    source = Path(args.source)
    if args.list:
        marks, total, why = scan_pptx(source)
        if why:
            print(why)
            return 1
        print(f"{total} 页，找到 {len(marks)} 个跨页固定对象：")
        for mark in marks:
            flag = "✓" if mark.eligible else "·"
            print(f"  {flag} [{mark.kind}] {mark.label}  {mark.hits}/{mark.total} 页 "
                  f"{mark.home} {mark.area_percent:.2f}%"
                  + (f"  （不动：{mark.reason}）" if mark.reason else ""))
            print(f"      key={mark.key}")
        return 0

    destination = Path(args.destination) if args.destination else \
        source.with_name(f"{source.stem}_已去水印.pptx")
    result = clean_pptx(source, destination, select=args.select)
    print(result.message)
    return 0 if result.success else 1


if __name__ == "__main__":
    raise SystemExit(main())
