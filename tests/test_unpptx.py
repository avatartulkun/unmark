"""PPTX 去角标的自包含测试。

跑法：cd 到本工具目录，执行  python3 -m pytest -q
"""

import copy
import hashlib
import io
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from unpptx import PptxOptions, clean_pptx, scan_pptx

WIDE, TALL = Inches(13.333), Inches(7.5)


def _deck(path: Path, slides: int = 6, badge: str = "Gemini Notebook",
          where: str = "slide", centered_repeat: bool = False) -> Path:
    """造一份「正文逐页在变、角标每页同一位置」的 PPTX。

    `where` 决定角标画在哪儿：slide（逐页放置）或 master（母版，只画一次
    但每页都会显示出来）。
    """
    presentation = Presentation()
    presentation.slide_width, presentation.slide_height = WIDE, TALL
    blank = presentation.slide_layouts[6]

    for index in range(slides):
        slide = presentation.slides.add_slide(blank)
        body = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        body.text_frame.text = f"第 {index + 1} 页正文 " + "x" * index
        if centered_repeat:                    # 每页都一样、但压在正文区中央
            middle = slide.shapes.add_textbox(Inches(4), Inches(3.4), Inches(5), Inches(0.6))
            middle.text_frame.text = "公司内部资料"
        if badge and where == "slide":
            box = slide.shapes.add_textbox(Inches(10.4), Inches(6.8), Inches(2.6), Inches(0.5))
            box.text_frame.text = badge
            box.text_frame.paragraphs[0].runs[0].font.size = Pt(12)

    if badge and where == "master":
        # python-pptx 的 MasterShapes 没有 add_textbox，照实际文件的做法：
        # 在页面上建好再把 XML 搬进母版。
        host = presentation.slides[0]
        box = host.shapes.add_textbox(Inches(10.4), Inches(6.8), Inches(2.6), Inches(0.5))
        box.text_frame.text = badge
        presentation.slide_masters[0].shapes._spTree.append(copy.deepcopy(box._element))
        box._element.getparent().remove(box._element)

    presentation.save(str(path))
    return path


def _texts(path: Path) -> list[str]:
    presentation = Presentation(str(path))
    out = [sh.text_frame.text for s in presentation.slides
           for sh in s.shapes if sh.has_text_frame]
    out += [sh.text_frame.text for m in presentation.slide_masters
            for sh in m.shapes if sh.has_text_frame]
    return out


def test_badge_placed_on_every_slide_is_removed(tmp_path: Path) -> None:
    """跨页同一位置、同一内容的角标要被删掉，正文一个字不动。"""
    source = _deck(tmp_path / "in.pptx")
    result = clean_pptx(source, tmp_path / "out.pptx")

    assert result.success, result.message
    assert result.removed == 6, "每页都该删掉一处"
    texts = _texts(tmp_path / "out.pptx")
    assert not any("Gemini" in t for t in texts), "角标还在"
    assert sum("页正文" in t for t in texts) == 6, "正文被误删了"


def test_source_file_is_never_touched(tmp_path: Path) -> None:
    """原文件永不改写——和 PDF 那条路同一条底线。"""
    source = _deck(tmp_path / "in.pptx")
    before = hashlib.sha1(source.read_bytes()).hexdigest()
    clean_pptx(source, tmp_path / "out.pptx")
    assert hashlib.sha1(source.read_bytes()).hexdigest() == before


def test_badge_on_the_master_is_found_and_removed_once(tmp_path: Path) -> None:
    """画在母版上的角标同样算数，而且只用删一次。

    母版上的一个形状会显示在所有用它的页面上，所以按「实际覆盖到多少页」算证据，
    和逐页放置的角标同一个标准；删在母版上则是一次性解决。
    """
    source = _deck(tmp_path / "in.pptx", badge="Made with NotebookLM", where="master")
    marks, total, why = scan_pptx(source)
    assert not why and total == 6
    badge = next(m for m in marks if "NotebookLM" in m.label)
    assert badge.home == "master"
    assert badge.hits == 6, "母版上的形状要按它覆盖到的页数计票"

    result = clean_pptx(source, tmp_path / "out.pptx")
    assert result.success, result.message
    assert result.removed == 1, "母版上删一处就够，不该逐页删"
    assert not any("NotebookLM" in t for t in _texts(tmp_path / "out.pptx"))


def test_picture_badge_is_matched_by_pixels(tmp_path: Path) -> None:
    """图片角标按像素内容认，改名换 id 都躲不掉。"""
    from PIL import Image

    buffer = io.BytesIO()
    Image.new("RGB", (160, 40), (240, 240, 240)).save(buffer, "PNG")
    presentation = Presentation()
    presentation.slide_width, presentation.slide_height = WIDE, TALL
    for index in range(6):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1)) \
            .text_frame.text = f"内容 {index}"
        buffer.seek(0)
        slide.shapes.add_picture(buffer, Inches(11.0), Inches(6.8), Inches(1.7), Inches(0.42))
    source = tmp_path / "in.pptx"
    presentation.save(str(source))

    marks, _total, _why = scan_pptx(source)
    assert any(m.kind == "picture" and m.eligible for m in marks)
    result = clean_pptx(source, tmp_path / "out.pptx")
    assert result.success, result.message
    out = Presentation(str(tmp_path / "out.pptx"))
    assert not any(sh.shape_type == 13 for s in out.slides for sh in s.shapes), "图片角标还在"


def test_repeated_block_inside_the_body_is_left_alone(tmp_path: Path) -> None:
    """压在正文区里的重复块不自动删。

    「每页都有」只说明它固定，不说明它是水印——页脚、栏目名、装饰块都满足。
    水印待在页边，正文待在中间，所以自动模式只动完全落在正文区之外的形状；
    压在中间的照样列出来，由人自己决定。
    """
    source = _deck(tmp_path / "in.pptx", badge="", centered_repeat=True)
    marks, _total, _why = scan_pptx(source)
    middle = next(m for m in marks if m.label == "公司内部资料")
    assert not middle.eligible
    assert "正文区" in middle.reason

    result = clean_pptx(source, tmp_path / "out.pptx")
    assert not result.success, "不该自动删正文区里的东西"
    assert "公司内部资料" in "".join(_texts(source))


def test_user_can_select_a_declined_candidate(tmp_path: Path) -> None:
    """自动模式不碰的，人可以点名要删。"""
    source = _deck(tmp_path / "in.pptx", badge="", centered_repeat=True)
    marks, _total, _why = scan_pptx(source)
    middle = next(m for m in marks if m.label == "公司内部资料")

    result = clean_pptx(source, tmp_path / "out.pptx", select=[middle.key])
    assert result.success, result.message
    assert not any("公司内部资料" in t for t in _texts(tmp_path / "out.pptx"))
    assert sum("页正文" in t for t in _texts(tmp_path / "out.pptx")) == 6


def test_unknown_selection_is_refused(tmp_path: Path) -> None:
    """点名一个不存在的对象要报错，不能悄悄什么都不做。"""
    source = _deck(tmp_path / "in.pptx")
    result = clean_pptx(source, tmp_path / "out.pptx", select=["text:deadbeef:1:2"])
    assert not result.success and "不在候选里" in result.message


def test_short_decks_are_not_judged_automatically(tmp_path: Path) -> None:
    """页数太少就别自动下判断——两页之间「都有」说明不了什么。"""
    source = _deck(tmp_path / "in.pptx", slides=2)
    marks, total, _why = scan_pptx(source)
    assert total == 2
    assert all(not m.eligible for m in marks)
    assert not clean_pptx(source, tmp_path / "out.pptx").success


def test_signature_ignores_per_slide_object_ids(tmp_path: Path) -> None:
    """指纹必须抹掉逐页递增的 id 和自动生成的 name。

    没文字也没图片的形状（纯色块、线条）只能拿自己的 XML 当指纹，而 PPTX 里
    每个形状的 `id` 是全篇递增的、`name` 常带序号。不抹掉的话同一个角标在每页
    的指纹都不一样，跨页交集永远为空——这条路会静默地什么都找不到。
    """
    from pptx.enum.shapes import MSO_SHAPE

    presentation = Presentation()
    presentation.slide_width, presentation.slide_height = WIDE, TALL
    for index in range(6):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1)) \
            .text_frame.text = f"内容 {index}"
        slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(11.2), Inches(6.9), Inches(1.6), Inches(0.4))
    source = tmp_path / "in.pptx"
    presentation.save(str(source))

    marks, _total, _why = scan_pptx(source)
    assert any(m.kind == "shape" and m.hits == 6 for m in marks), \
        "纯图形角标没被跨页认出来（多半是 id/name 没抹干净）"


def test_verification_catches_a_bad_delete(tmp_path: Path, monkeypatch) -> None:
    """复核必须拦得住「删过头」。

    删除很容易做过头，而且不像涂像素那样一眼看得出来，所以不信内存里的状态，
    只信重新读回来的文件。
    """
    import unpptx

    source = _deck(tmp_path / "in.pptx")
    real_drop = unpptx._drop
    state = {"extra": True}

    def greedy(shape):
        """删角标的同时，顺手把它所在页的正文也删掉——复核应当发现。"""
        parent = shape._element.getparent()
        real_drop(shape)
        if state["extra"] and parent is not None and len(parent) > 1:
            parent.remove(parent[-1])
            state["extra"] = False

    monkeypatch.setattr(unpptx, "_drop", greedy)
    result = clean_pptx(source, tmp_path / "out.pptx")
    assert not result.success and "已中止" in result.message
    assert not (tmp_path / "out.pptx").exists(), "复核不过就不该留下输出文件"


def test_output_keeps_slide_count_and_shrinks(tmp_path: Path) -> None:
    """删对象是无损操作：页数不变，体积只减不增。

    和 PDF 那条路正相反——那边要把位图无损重编码，输出常常涨到几倍。
    """
    source = _deck(tmp_path / "in.pptx")
    destination = tmp_path / "out.pptx"
    assert clean_pptx(source, destination).success
    assert len(Presentation(str(destination)).slides._sldIdLst) == 6
    assert destination.stat().st_size <= source.stat().st_size


def test_scan_reports_why_a_candidate_was_declined(tmp_path: Path) -> None:
    """列表要把「为什么不动它」说出来，不能只给一个空结果。"""
    source = _deck(tmp_path / "in.pptx", badge="", centered_repeat=True)
    result = clean_pptx(source, tmp_path / "out.pptx")
    assert not result.success
    assert "公司内部资料" in result.message and "正文区" in result.message


def test_outer_only_can_be_turned_off(tmp_path: Path) -> None:
    """位置这道闸门是可以关的——关掉就只剩「跨页固定 + 够小」。"""
    source = _deck(tmp_path / "in.pptx", badge="", centered_repeat=True)
    marks, _total, _why = scan_pptx(source, PptxOptions(outer_only=False))
    assert any(m.label == "公司内部资料" and m.eligible for m in marks)
